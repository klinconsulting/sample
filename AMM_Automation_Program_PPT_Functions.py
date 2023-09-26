from pptx.chart.data import CategoryChartData,XySeriesData, XyChartData
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.dml import MSO_THEME_COLOR
from PIL import Image
from pptx.chart.plot import BarPlot
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LABEL_POSITION


def replace_text_in_slide(presentation, slidenumber, shape_name, placeholder, new_value,\
                          set_format=False,set_color=False ,font='Amazon Ember Display', font_size=36, alignment='CENTER', superscript_text='',font_color=[128,130,133], bold=False):
    '''
        Description:
        This function replace a set of characters with a value provided by the user.
        The replacement is done only in the shape and slide specified by the user.

    Parameters:
        :param presentation: object
            The presentation objected created with python-pptx. Usually it is called: ppt.
        :param slidenumber:  integer
            The number of the slide as reported in the template.
        :param shape_name: string
            The name of the shape you want to change in the slide.
            You can get this name from the selection panel. Home -> Arrange -> Selection panel
        :param placeholder: string
            The characters that are the placeholder for the new value: example "XX" or "<Brand_Name>"
        :param new_value: variable/str/int
            The new value you want to insert in the PPT.
        :param set_format: boolean
            This variable indicates whether we want to set format for the text. Sometimes, replace function destroys the 
            original text format. The default value is False.
        :param font: string
            The string represents the font style. The default value is Amazon Ember Display.
        :param font_size: integer 
            The integer indicates the font size. The default value is 36.
        :param alignment: string
            The string represents on the alignment of text on the slide
        :param superscript_text: string
            It represents the superscript text to be added at the end of the paragragh.
        :param bold: boolean
            Boolean indicator on whether the font should be bold
    Return:
        Nothing
    '''
    
    slide = presentation.slides[slidenumber-1]
    for shape in slide.shapes:
        if shape.name == shape_name:
            shape.text = shape.text.replace(str(placeholder), str(new_value))
            
            if set_format:
                from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
                alignment_dictonary = {
                'CENTER': PP_PARAGRAPH_ALIGNMENT.CENTER,
                'LEFT': PP_PARAGRAPH_ALIGNMENT.LEFT,
                'RIGHT': PP_PARAGRAPH_ALIGNMENT.RIGHT,
                'JUSTIFY': PP_PARAGRAPH_ALIGNMENT.JUSTIFY}
                
                for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.size = Pt(font_size)
                        paragraph.font.name = font
                        paragraph.font.bold = bold
                        paragraph.alignment = alignment_dictonary[alignment]
                        
                        if len(superscript_text) > 0:
                            run = paragraph.add_run()
                            run.text=superscript_text
                            set_superscript(run.font)
                            
                        if set_color:
                            paragraph.font.color.rgb = RGBColor(font_color[0], font_color[1], font_color[2])

def replace_text_in_slide_v2(presentation, slidenumber, shape_name,
                                       placeholder='',new_value='', tag='higher',
                                       font_style='Amazon Ember Light', font_size=12, bold=False, alignment='CENTER',
                                       font_red=120, font_green=120, font_blue=120,using_font_color=False,
                                       color_code = MSO_THEME_COLOR.BACKGROUND_2,
                                       symbol='x ' 
                            ):
    '''
        Description:
        This function replace a set of characters with a value provided by the user.
        The replacement is done only in the shape and slide specified by the user.

    Parameters:
        :param presentation: object
            The presentation objected created with python-pptx. Usually it is called: ppt.
        :param slidenumber:  integer
            The number of the slide as reported in the template.
        :param shape_name: string
            The name of the shape you want to change in the slide.
            You can get this name from the selection panel. Home -> Arrange -> Selection panel
        :param placeholder: string
            The characters that are the placeholder for the new value: example "XX" or "<Brand_Name>"
        :param new_value: variable/str/int
            The new value you want to insert in the PPT.
        :param tag: string
            Text to be inserted into the PPT
        :param font_style: string
            Text font style
        :param font_size: int
            Text size
        :param bold: boolean
            Text to be bold or not
        :param alignment: string
            Text alignment
        :param font_red: int
            Numerical value for red in the RGB code
        :param font_green: int
            Numerical value for green in the RGB code
        :param font_blue: int
            Numerical value for blue in the RGB code
        :param using_font_color: boolean
            Text font using RGB code or not
        :param color_code: string
            color theme
        :param symbol: string
            string for symbol such as 'x' if any
        
    Return:
        Nothing
    '''
    from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
    alignment_dictonary = {
        'CENTER': PP_PARAGRAPH_ALIGNMENT.CENTER,
        'LEFT': PP_PARAGRAPH_ALIGNMENT.LEFT,
        'RIGHT': PP_PARAGRAPH_ALIGNMENT.RIGHT,
        'JUSTIFY': PP_PARAGRAPH_ALIGNMENT.JUSTIFY}
    
    slide = presentation.slides[slidenumber-1]
    
    for shape in slide.shapes:
        if shape.name == shape_name: 
            text_frame = shape.text_frame
            text_frame.clear()
            
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(new_value) + symbol + tag
            p.alignment = alignment_dictonary[alignment]
            
            font = run.font
            font.name = font_style
            font.size = Pt(font_size)
            font.bold = bold
            
            
            if using_font_color:
                font.color.rgb = RGBColor(font_red, font_green, font_blue)
            else:
                font.color.theme_color = color_code
            
            
def replace_text_in_slide_v3(presentation, slidenumber, shape_name,
                                       placeholder='',new_value='', tag='higher', tag_size=26, tag_color=[120,120,120],
                                       font_style='Amazon Ember Display', font_size=12, bold=False, alignment='CENTER',
                                       font_red=120, font_green=120, font_blue=120,
                                       color_code = MSO_THEME_COLOR.BACKGROUND_2,
                                       symbol='x ', symbol_size=26, symbol_color=[120,120,120]
                            ):
    '''
        Description:
        This function replace a set of characters with a value provided by the user.
        The replacement is done only in the shape and slide specified by the user.

    Parameters:
        :param presentation: object
            The presentation objected created with python-pptx. Usually it is called: ppt.
        :param slidenumber:  integer
            The number of the slide as reported in the template.
        :param shape_name: string
            The name of the shape you want to change in the slide.
            You can get this name from the selection panel. Home -> Arrange -> Selection panel
        :param placeholder: string
            The characters that are the placeholder for the new value: example "XX" or "<Brand_Name>"
        :param new_value: variable/str/int
            The new value you want to insert in the PPT.
        :param tag: string
            Text to be inserted into the PPT
        :param tag_size: int
            Text size for tag
        :param tag_color: list of int
            RGB code for tag
        :param font_style: string
            Text font style
        :param font_size: int
            Text size
        :param bold: boolean
            Text to be bold or not
        :param alignment: string
            Text alignment
        :param font_red: int
            Numerical value for red in the RGB code
        :param font_green: int
            Numerical value for green in the RGB code
        :param font_blue: int
            Numerical value for blue in the RGB code
        :param color_code: string
            color theme
        :param symbol: string
            string for symbol such as 'x' if any
        :param symbol_size: int
            Text size for symbol
        :param symbol_color: list of int
            RGB code for symbol
    Return:
        Nothing
    '''
    from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
    alignment_dictonary = {
        'CENTER': PP_PARAGRAPH_ALIGNMENT.CENTER,
        'LEFT': PP_PARAGRAPH_ALIGNMENT.LEFT,
        'RIGHT': PP_PARAGRAPH_ALIGNMENT.RIGHT,
        'JUSTIFY': PP_PARAGRAPH_ALIGNMENT.JUSTIFY}
    
    slide = presentation.slides[slidenumber-1]
    
    for shape in slide.shapes:
        if shape.name == shape_name: 
            text_frame = shape.text_frame
            text_frame.clear()
            
            # new value
            p = text_frame.add_paragraph()
            #p.space_before = Pt(0)
            p.line_spacing = Pt(0)
            run = p.add_run()
            run.text = str(new_value)
            p.alignment = alignment_dictonary[alignment]
            
            font = run.font
            font.name = font_style
            font.size = Pt(font_size)
            font.bold = bold
            font.color.theme_color = color_code
            
            
            # symbol
            p_symbol = text_frame.add_paragraph()
            #p_symbol.add_line_break()
            p_symbol.space_before = Pt(20)
            #p_symbol.line_spacing = Pt(100)
            run_symbol = p_symbol.add_run()
            run_symbol.text = symbol
            p_symbol.alignment = alignment_dictonary[alignment]
            
            font_symbol = run_symbol.font
            font_symbol.name = font_style
            font_symbol.size = Pt(symbol_size)
            font_symbol.bold = bold
            font_symbol.color.rgb = RGBColor(symbol_color[0],symbol_color[1],symbol_color[2])
            
            # tag
            
            p_tag = text_frame.add_paragraph()
            run_tag = p_symbol.add_run()
            run_tag.text = tag
            p_tag.alignment = alignment_dictonary[alignment]
            
            font_tag = run_tag.font
            font_tag.name = font_style
            font_tag.size = Pt(tag_size)
            font_tag.bold = bold
            font_tag.color.rgb = RGBColor(tag_color[0], tag_color[1], tag_color[2])
            
def replace_text_in_notes(presentation, slidenumber, placeholder, new_value):
    '''
        Description:
        This function replace a set of characters with a value provided by the user in the notes of the slide specified by the user.

    Parameters:
        :param presentation: object
            The presentation objected created with python-pptx. Usually it is called: ppt.
        :param slidenumber:  integer
            The number of the slide as reported in the template.
        :param placeholder: string
            The characters that are the placeholder for the new value: example "XX" or "<Brand_Name>"
        :param new_value: variable/str/int
            The new value you want to insert in the PPT.

    Return:
        Nothing
    '''
    slide = presentation.slides[slidenumber-1]
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = text_frame.text.replace(placeholder, str(new_value))     

def replace_text_in_group(presentation, slidenumber,group_name,shape_name, placeholder, new_value, font='Amazon Ember Display', font_size=26, bold=False, superscript_text='', alignment='CENTER'):
    
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    
    slide = presentation.slides[slidenumber-1]
    group_shapes = [
        shp for shp in slide.shapes
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP
    ]
    
    for group_shape in group_shapes:
        if group_shape.name == group_name:
            for shape in group_shape.shapes:
                if shape.name == shape_name:
                    shape.text = shape.text.replace(str(placeholder), str(new_value))
                
                    for paragraph in shape.text_frame.paragraphs:
               
                        paragraph.font.size = Pt(font_size)
                        paragraph.font.name = font
                        paragraph.font.color.rgb = RGBColor(120,120,120)
            
                        if len(superscript_text) > 0:
                            run = paragraph.add_run()
                            run.text=superscript_text
                            set_superscript(run.font)
            
def add_image_in_slide(presentation, slidenumber, shape_name, file_name):
    '''
        Description:
        This function add an image to a PPT Picture Placeholder.
        The slide need to contain a Picture Placeholder inserted from the Master Slide View.
        The image need to be saved in the same folder of the jupyter notebook running.

    Parameters:
        :param presentation: object
            The presentation objected created with python-pptx. Usually it is called: ppt.
        :param slidenumber:  integer
            The number of the slide as reported in the template.
        :param shape_name: string
            The name of the shape you want to change in the slide.
            You can get this name from the selection panel. Home -> Arrange -> Selection panel
            The shape need to be a Picture Placeholder.
        :param file_name: string
            The name of the picture.
            For example, if it is a picture of an ASIN, that you have dowload from a dataframe:
            df['asin'].to_list()[0] + '.jpg'
            if it is a picture of a generate chart it could look like:
            'Chart_Landscale_' + category + '.pgn'

    Return:
        Nothing
    '''
    from io import FileIO
    slide = presentation.slides[slidenumber-1]
    for shape in slide.shapes:
        if shape.name == shape_name:
            shape.insert_picture(FileIO(file_name))

def replace_image_in_slide(presentation, slidenumber, file_name, left_spec=1,top_spec=1, width_spec=3, height_spec=4, image_type = 'PICTURE'):
    '''
        Description:
        This function add an image to a PPT Picture Placeholder.
        The slide need to contain a Picture Placeholder inserted from the Master Slide View.
        The image need to be saved in the same folder of the jupyter notebook running.

    Parameters:
        :param presentation: object
            The presentation objected created with python-pptx. Usually it is called: ppt.
        :param slidenumber:  integer
            The number of the slide as reported in the template.
        :param file_name: string
            The name of the picture.
            For example, if it is a picture of an ASIN, that you have dowload from a dataframe:
            df['asin'].to_list()[0] + '.jpg'
            if it is a picture of a generate chart it could look like:
            'Chart_Landscale_' + category + '.pgn'
        :param left_sepc: float
            It represents the left margin on the slide
        :param top_spec: float
            It represents the top margin on the slide
        :param width_spec: float
            It represents the width of the image on the slide
        :param height_spec: float
            It represents the height of the image on the slide
        :param image_type: string
            It represents the element type on the slide

    Return:
        Nothing
    '''
    from io import FileIO
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE

    left = Inches(left_spec)
    top = Inches(top_spec)
    width = Inches(width_spec)
    height = Inches(height_spec)
    
    slide = presentation.slides[slidenumber-1]
    old_shape = None
    # remove old picture
    for shape in slide.shapes:
        if str(shape.shape_type).find(image_type)!=-1:
            # found the picture, assign the shape to old element and remove it later
            old_shape = shape
    
    new_shape=slide.shapes.add_picture(FileIO(file_name),left,top,width,height)

    if old_shape is not None:
        old_shape = old_shape._element
        new_shape = new_shape._element
        old_shape.addnext(new_shape)  # 
        old_shape.getparent().remove(old_shape)

    # adjust new shape specification
            
def change_bar_chart(presentation, slidenumber, chart_name, series_names, df, x_column, y_columns,\
                     ChartType = 'Bar', set_format=False, marker_size=22):
    '''
        Description:
        This function replace the data in a chart within the PPT with the data you want to show.

        This functions works for:
            - BAR CHARTS (verical and horizontal / stack or unstaked)
            - LINE CHARTS.

        This function won't work with COMBO CHARTS.

        First make sure to create a template chart in the PPT with dummy numbers already formatted as you want.
        Then use this function to replace the data. If you have more than one series, pay attention to the order.
        For example, if your chart has 2 series, 1 blue and 1 grey in this order,
        if you pass to series ['Sales', 'Sales Year Ago'], sales will become the blue (first) and sales y.a. grey.

    Parameters:
        :param presentation: object
            The presentation objected created with python-pptx. Usually it is called: ppt.
        :param slidenumber:  integer
            The number of the slide as reported in the template.
        :param chart_name: string
            The name of the shape that contains the chart.
            You can get this name from the selection panel. Home -> Arrange -> Selection panel
        :param series_names: list of string
            The name of the series in the final chart as you want to be displayed in the legend.
            Pay attention that the order machtes the order of the series in your template chart.
            Pass some final names, as string in a list:
            Example: ['Percentage of new to brand customers', 'Percentage of repeat purchasers']
            Note: even if you have only 1 serie like 'Percentage of new to brand customers'
            you still need to put within list like:
            ['Percentage of new to brand customers']
        :param df: Pandas DataFrame
            The DataFrame that contains the columns you want to use to replace.
        :param x_column: string
            Name of the column that contains the value to insert in the X axis
        :param y_columns: list of string
            List of the names of the columns that contains the value to insert in the Y axis
            Make sure to match the order of the series name.
            Example: ['percentage_NTB', 'percentage_repeat']
            Note: even if you have only 1 serie like 'percentage_NTB' you still need to put within list like:
            ['percentage_NTB']
        :param ChartType: string indicating type of chart
            Note: this function is applicable to a few chart types even the original name is "change_bar_chart".
        :param set_format: boolean
            Boolean indicator of whether we should set font size, color, bar width and font
        : param marker_size: integer
            marker size of the line chart

    Return:
        Nothing
    '''

    # prepare the data
    chart_data = CategoryChartData()

    # change the X values
    
    chart_data.categories = df[x_column].to_list()

    # create the tuple of the data
    for serie_name, y_name in zip(series_names, y_columns):
        t = ()
        for value in df[y_name].values:
            t = t + (value,)
        # add the series of each Y
        chart_data.add_series(serie_name, t)

    # substitute the data in the chart
    slide = presentation.slides[slidenumber-1]
    for shape in slide.shapes:
        if shape.name == chart_name:
            chart = shape.chart
            chart.replace_data(chart_data)
    
    if set_format:
        
        if str(chart.chart_type).find('COLUMN_CLUSTERED')!=-1:
            
            # change the width of the bar
            chart.x = 100
            
        if str(chart.chart_type).find('LINE_MARKER')!=-1:
            
             # dark blue: RGB(62,65,67), blue-grey: RGB(102,153,204)
             # gray, text 2 (128,128,128), orange: RGB(255,90,0) 
            
            fill = chart.series[0].marker.format.fill
            fill.solid()
            
            # 35,47,63 -> DARK BLUE, 166, 166, 166 -> GREY
            fill.fore_color.rgb = RGBColor(166,166,166)
            chart.series[0].marker.size = marker_size
            
            chart.series[0].format.line.width = Pt(3)
                   
def dataframe_to_table(presentation, slidenumber, table_name,
                       dataframe,
                       font='Amazon Ember Light', font_size=12, bold=False,
                       alignment='CENTER',
                       font_red=120, font_green=120, font_blue=120,
                       cell_red=255, cell_green=255, cell_blue=255):
    '''
        Description:
            This function will add the data from a dataframe and place it in a table in the PPT.

            NOTE 1: This function won't copy the dataframe columns names. It will copy only the values.

            NOTE 2: Every time you change a table, it turns back to the awful power point detfault.
            For this reason every time you change a table you need to pass all the formatting or using
            the default of the function.

            NOTE 3: If you don't specify the layout parameters above, they will take their default value:
            font='Amazon Ember Light', font_size=12, bold=False, alignment='CENTER',
            font colur = grey, cell background = white.

            NOTE 4: Make sure your dataframe is formatted like you want the PPT to look. For example,
            if you want to create a percentage column, you need to build a column
            in the dataframe like '98%' as string.

        Parameters:
        :param presentation: object
            The presentation objected created with python-pptx. Usually it is called: ppt.
        :param slidenumber:  integer
            The number of the slide as reported in the template.
        :param table_name: string
            The name of the shape that contains the table.
            You can get this name from the selection panel. Home -> Arrange -> Selection panel
        :param dataframe: Pandas DataFrame
            The dataframe where the data is saved.

        Optional Parameters:
        :param font: string
            Enter the font you want to use in the table. Example: Amazon Ember.
        :param font_size: integer
            Enter the font size. Example: 12
        :param bold: boolean (True/False)
            This specify if you want the font to be bold or not:
            True -> Bold; False -> Not bold

        :param font_red: integer - range 0-255
        :param font_green: integer - range 0-255
        :param font_blue: integer - range 0-255
            These three integers will set up the RGB color of the font in the table.
            For example: font_red=0,  font_green=90, font_blue=149 is Amazon Blue.

        :param cell_red: integer - range 0-255
        :param cell_green: integer - range 0-255
        :param cell_blue: integer - range 0-255
            These three integers will set up the RGB color of the background of each cell in the table.
            For example: font_red=255,  font_green=255, font_blue=255 is white background.

        :param alignment: string
            This parameter can take only the following value:
            CENTER, LEFT, RIGHT, JUSTIFY
            They command the allignment of the text in the table.

        Return:
            Nothing
        '''

    from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
    alignment_dictonary = {
        'CENTER': PP_PARAGRAPH_ALIGNMENT.CENTER,
        'LEFT': PP_PARAGRAPH_ALIGNMENT.LEFT,
        'RIGHT': PP_PARAGRAPH_ALIGNMENT.RIGHT,
        'JUSTIFY': PP_PARAGRAPH_ALIGNMENT.JUSTIFY}

    slide = presentation.slides[slidenumber - 1]
    for shape in slide.shapes:
        if shape.name == table_name:
            table = shape.table
            for row in range(dataframe.shape[0]):
                for col in range(dataframe.shape[1]):
                    # cell start with 0, 0 as top left column
                    # starting at row+1 we start from the first row in the PPT table
                    cell = table.cell(row+1, col)  # row, columns
                    cell.text = str(dataframe.values[row][col])
                    cell.text_frame.paragraphs[0].alignment = alignment_dictonary[alignment]
                    
                    table.cell(row+1, col).text_frame.paragraphs[0].font.size = Pt(font_size)
                    table.cell(row+1, col).text_frame.paragraphs[0].font.name = font
                    table.cell(row+1, col).text_frame.paragraphs[0].font.bold = bold
                    table.cell(row+1, col).text_frame.paragraphs[0].font.color.rgb = RGBColor(font_red, font_green,
                                                                                            font_blue)
                    table.cell(row+1, col).fill.solid()
                    table.cell(row+1, col).fill.fore_color.rgb = RGBColor(cell_red, cell_green, cell_blue)




def add_data_to_specific_cell_in_table(presentation, slidenumber, table_name,
                                       row, col, content,
                                       font='Amazon Ember Light', font_size=12, bold=False, alignment='CENTER',
                                       font_red=120, font_green=120, font_blue=120,
                                       cell_red=255, cell_green=255, cell_blue=255
                           ):
    '''
        Description:
            This function will write a specific value from a dataframe
            and place it in a specific cell in the table in the PPT.

            NOTE 1: Every time you change a table, it turns back to the awful power point detfault.
            For this reason every time you change a table you need to pass all the formatting or using
            the default of the function.

            NOTE 2: If you don't specify the layout parameters above, they will take their default value:
            font='Amazon Ember Light', font_size=12, bold=False, alignment='CENTER',
            font colur = grey, cell background = white.

        Parameters:
        :param presentation: object
            The presentation objected created with python-pptx. Usually it is called: ppt.

        :param slidenumber:  integer
            The number of the slide as reported in the template.

        :param table_name: string
            The name of the shape that contains the table.
            You can get this name from the selection panel. Home -> Arrange -> Selection panel

        :param row: integer
            The row number of the table. The header (first row from the top) is 1.
            The second row is 2 and so on.

        :param col: integer
            The column number of the table. The first column on the left is number 1.
            The second column is then 2 and so on.

        :param content: string/integer
            The value that you want to enter in the PPT.

        Optional Parameters:
        :param font: string
            Enter the font you want to use in the table. Example: Amazon Ember.
        :param font_size: integer
            Enter the font size. Example: 12
        :param bold: boolean (True/False)
            This specify if you want the font to be bold or not:
            True -> Bold; False -> Not bold

        :param font_red: integer - range 0-255
        :param font_green: integer - range 0-255
        :param font_blue: integer - range 0-255
            These three integers will set up the RGB color of the font in the table.
            For example: font_red=0,  font_green=90, font_blue=149 is Amazon Blue.

        :param cell_red: integer - range 0-255
        :param cell_green: integer - range 0-255
        :param cell_blue: integer - range 0-255
            These three integers will set up the RGB color of the background of each cell in the table.
            For example: font_red=255,  font_green=255, font_blue=255 is white background.

        :param alignment: string
            This parameter can take only the following value:
            CENTER, LEFT, RIGHT, JUSTIFY
            They command the allignment of the text in the table.

        Return:
            Nothing
        '''

    # NOTE: row and columns start from zero!
    # bold is true/false
    from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
    alignment_dictonary = {
        'CENTER': PP_PARAGRAPH_ALIGNMENT.CENTER,
        'LEFT': PP_PARAGRAPH_ALIGNMENT.LEFT,
        'RIGHT': PP_PARAGRAPH_ALIGNMENT.RIGHT,
        'JUSTIFY': PP_PARAGRAPH_ALIGNMENT.JUSTIFY}

    slide = presentation.slides[slidenumber - 1]
    for shape in slide.shapes:
        if shape.name == table_name:
            table = shape.table
            cell = table.cell(row-1, col-1)  # row, columns
            cell.text = content
            cell.text_frame.paragraphs[0].alignment = alignment_dictonary[alignment]
            table.cell(row-1, col-1).text_frame.paragraphs[0].font.size = Pt(font_size)
            table.cell(row-1, col-1).text_frame.paragraphs[0].font.name = font
            table.cell(row-1, col-1).text_frame.paragraphs[0].font.bold = bold
            table.cell(row-1, col-1).text_frame.paragraphs[0].font.color.rgb = RGBColor(font_red, font_green, font_blue)
            table.cell(row-1, col-1).fill.solid()
            table.cell(row-1, col-1).fill.fore_color.rgb = RGBColor(cell_red, cell_green, cell_blue)


def change_chart_data_labels(presentation, slidenumber,
                             chart_name, chart_serie_name,
                             df, label_column, set_format=False, font='Amazon Ember Display',\
                             font_size=36, bold=True, ChartType='Bar',Numeric=False):
    '''
        Description:
            This function will take a column of a dataframe and add to a chart as custom text labels.
            This function should be used when you want to add labels different from the data.

        Parameters:
            :param presentation: object
                The presentation objected created with python-pptx. Usually it is called: ppt.

            :param slidenumber:  integer
                The number of the slide as reported in the template.

            :param chart_name: string
                The name of the shape that contains the chart.
                You can get this name from the selection panel. Home -> Arrange -> Selection panel

            :param chart_serie_name: string
                The name of the series that you want to add data labels to.

            :param df: Pandas DataFrame
                The DataFrame that contains the columns you want to use as data label.

            :param label_column: string
                Name of the column that contains the value to insert as data label to the chart series.
            :param set_format: boolean
                Boolean indicator of whether to set font, font size, font color and label alignment
            :param font: string
                Font style
            :param font_size: integer
                Font size
            :param bold: boolean
                Boolean indicator of whether font should be bold
            :param ChartType: string
                Chart type
            :param Numeric: boolean
                Boolean indicator of whether input label column is numeric
        Return:
            Nothing
    '''

    #Change the lables only after you changed the underline data.
    list_of_lables = df[label_column].to_list()
    
    if Numeric:
        list_of_lables = [round(float(num),1) for num in list_of_lables]
        
    slide = presentation.slides[slidenumber-1]
    for shape in slide.shapes:
        if shape.name == chart_name:
            chart = shape.chart
            for series in chart.series:
                if series.name == chart_serie_name:
                    for k in range(len(series.points)):
                        series.points[k].data_label.text_frame.text = str(list_of_lables[k])
                        
                        if set_format:
                            series.points[k].data_label.text_frame.paragraphs[0].font.size = Pt(font_size)
                            #series.points[k].data_label.font.name = font
                            #series.points[k].data_label.font.bold = bold
                            series.points[k].data_label.text_frame.paragraphs[0].font.name = font 
                            
                            if ChartType != 'Bar':
                                # orange: RGB(62,65,67), blue-grey: RGB(102,153,204)
                                # gray, text 2 (128,128,128)
                                
                                series.points[k].data_label.position = XL_LABEL_POSITION.ABOVE
                                #series.points[k].data_label.font.color.rgb = RGBColor(255,0,0)
                                series.points[k].data_label.text_frame.paragraphs[0].font.color.rg=\
                                RGBColor(128,130,133)
                                
                                
                                # 128,128,128 -> dark blue; 35,47,63 ?  166,166,166 -> grey
                                # 127, 127,127
                                
                            else:
                                # 128,130,133 -> grey, 0,0,0 -> black
                                series.points[k].data_label.text_frame.paragraphs[0].font.color.rg=\
                                RGBColor(128,130,133)

def insert_error_shape(presentation, slidenumber, text_error):
    '''
        Description:
            This function will create a red rectangle with a custom text message.
            You should use this fucntion to aller the user when a slide need to modified or deleted.
            Use this function together with the and IF function to decide when and what to show.

        Parameters:
            :param presentation: object
                The presentation objected created with python-pptx. Usually it is called: ppt.

            :param slidenumber:  integer
                The number of the slide as reported in the template.

            :param text_error: string
                The text you want to display on the shape to inform the user.

        Return:
            Nothing
    '''

    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    left = Inches(1.0)
    top = Inches(0.0)
    width = Inches(8.0)
    height = Inches(2.0)

    shapes = presentation.slides[slidenumber - 1].shapes
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

    shape_fill = shape.fill
    shape_fill.solid()
    shape_fill.fore_color.rgb = RGBColor(255, 0, 0)
    shape_line = shape.line
    shape_line.color.rgb = RGBColor(255, 0, 0)

    shape.text = text_error


def apt_ranges(number):
    '''
        Description:
            This function take a number and return the equivalent Audiance Planning Tool (APT) range.

        Parameters:
            :param presentation: number
                The number you want to transform in ranges.

        Return:
            The APT Range as a string. Example: '1-5k' or '10-12M'.
    '''

    dict_ranges = {
        '1 - 5K': [0, 5000],
        '5K - 10K': [5000, 10000],
        '10K - 20K': [10000, 20000],
        '20K - 30K': [20000, 30000],
        '30K - 40K': [30000, 40000],
        '40K - 50K': [40000, 50000],
        '50K - 60K': [50000, 60000],
        '60K - 70K': [60000, 70000],
        '70K - 80K': [70000, 80000],
        '80K - 90K': [80000, 90000],
        '90K - 100K': [90000, 100000],
        '100K - 150K': [100000, 150000],
        '150K - 200K': [150000, 200000],
        '200K - 250K': [200000, 250000],
        '250K - 300K': [250000, 300000],
        '300K - 350K': [300000, 350000],
        '350K - 400K': [350000, 400000],
        '400K - 450K': [400000, 450000],
        '450K - 500K': [450000, 500000],
        '500K - 600K': [500000, 600000],
        '600K - 700K': [600000, 700000],
        '700K - 800K': [700000, 800000],
        '800K - 900K': [800000, 900000],
        '900K - 1M': [900000, 1000000],
        '1M - 1.5M': [1000000, 1500000],
        '1.5M - 2M': [1500000, 2000000],
        '2M - 2.5M': [2000000, 2500000],
        '2.5M - 3M': [2500000, 3000000],
        '3M - 3.5M': [3000000, 3500000],
        '3.5M - 4M': [3500000, 4000000],
        '4M - 4.5M': [4000000, 4500000],
        '4.5M - 5M': [4500000, 5000000],
        '5M - 6M': [5000000, 6000000],
        '6M - 7M': [6000000, 7000000],
        '7M - 8M': [7000000, 8000000],
        '8M - 9M': [8000000, 9000000],
        '9M - 10M': [9000000, 10000000],
        '10M - 12M': [10000000, 12000000],
        '12M - 14M': [12000000, 14000000],
        '14M - 16M': [14000000, 16000000],
        '16M - 18M': [16000000, 18000000],
        '18M - 20M': [18000000, 20000000],
        '20M - 25M': [20000000, 25000000],
        '25M - 30M': [25000000, 30000000],
        '30M - 35M': [30000000, 35000000],
        '35M - 40M': [35000000, 40000000],
        '40M - 45M': [40000000, 45000000],
        '45M - 50M': [45000000, 50000000],
        '50M - 60M': [50000000, 60000000],
        '60M - 70M': [60000000, 70000000],
        '70M - 80M': [70000000, 80000000],
        '80M - 90M': [80000000, 90000000],
        '90M - 100M': [90000000, 100000000]
    }
    for apt_range in dict_ranges.keys():
        if number >= dict_ranges[apt_range][0] and number < dict_ranges[apt_range][1]:
            return apt_range

    return '100M or more'


def Time_Series_Slope(data, days):
    import numpy as np
    coeffs = np.polyfit((data.index.values + 1)*days, list(data), deg=1)
    slope = coeffs[0]
    if slope > 0:
        return 'increasing'
    if slope == 0:
        return 'steady'
    if slope < 0:
        return 'decreasing'


def format_sales_units(number, currency=''):
    '''
        Description:
            This function take a number and a currency symbol and return a number formated for a PPT: example:
            25 --> €25
            235 --> €235
            2350 --> €2.4K
            23500 --> €23K
            235000 --> €235K
            2350000 --> €2.4MM
            23500000 --> €23MM
            235000000 --> €235MM

        Parameters:
            :param number: float or integer
                The number you want to transform.

            :param currency: string
                The currency symbol to add to the number. This value is optional

        Return:
            The number formatted.
    '''
    if number >= 10000000:  # 100MM
        number_formatted = currency + str('{:,}'.format(int(number / 1000000))) + 'MM'
        return number_formatted

    if number >= 1000000:  # 1.1MM
        number_formatted = currency + str('{:,}'.format(round(number / 1000000, 1))) + 'MM'
        return number_formatted

    if number >= 100000:  # 100k
        number_formatted = currency + str('{:,}'.format(int(number / 1000))) + 'K'
        return number_formatted

    if number >= 10000:  # 100k
        number_formatted = currency + str('{:,}'.format(int(number / 1000))) + 'K'
        return number_formatted

    if number >= 1000:  # 1.2k
        number_formatted = currency + str('{:,}'.format(round(number / 1000, 1))) + 'K'
        return number_formatted

    if number < 1000:  # 953
        number_formatted = currency + str(number)
        return number_formatted


def format_percentages(number, decimals):
    '''
        Description:
            This function take a number and transform into percentage (string). Example:
            0.234 -> 23.4%

        Parameters:
            :param number: float
                The number you want to transform.

            :param decimals: integer
                The currency symbol to add to the number. This value is optional

        Return:
            A string with the percentage value
    '''
    if decimals == 0:
        number_formatted = str(int(number * 100)) + '%'
        return number_formatted

    if decimals > 0:
        number_formatted = str(round(number * 100, decimals)) + '%'
        return number_formatted
    
    
def set_subscript(font):
    font._element.set('baseline', '-25000')

def set_superscript(font):
    font._element.set('baseline', '30000')

def set_strikethrough(font):
    font._element.set('strike','sngStrike')
    
def df_to_table_AMM_layout(presentation, slidenumber, table_name,
                            dataframe,
                            font='Amazon Ember Light', font_size=12, bold=False,
                            alignment='CENTER',
                            font_red=120, font_green=120, font_blue=120):

    from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
    alignment_dictonary = {
            'CENTER': PP_PARAGRAPH_ALIGNMENT.CENTER,
            'LEFT': PP_PARAGRAPH_ALIGNMENT.LEFT,
            'RIGHT': PP_PARAGRAPH_ALIGNMENT.RIGHT,
            'JUSTIFY': PP_PARAGRAPH_ALIGNMENT.JUSTIFY}
    slide = presentation.slides[slidenumber - 1]
    for shape in slide.shapes:
        if shape.name == table_name:
            table = shape.table
            for row in range(dataframe.shape[0]):
                if (row % 2) == 0:
                    cell_red = 242
                    cell_green = 242
                    cell_blue = 242

                else:
                    cell_red = 255
                    cell_green = 255
                    cell_blue = 255

                    for col in range(dataframe.shape[1]):
                        # cell start with 0, 0 as top left column
                        # starting at row+1 we start from the first row in the PPT table
                        cell = table.cell(row+1, col)  # row, columns
                        cell.text = str(dataframe.values[row][col])
                        cell.text_frame.paragraphs[0].alignment = alignment_dictonary[alignment]
                        table.cell(row+1, col).text_frame.paragraphs[0].font.size = Pt(font_size)
                        table.cell(row+1, col).text_frame.paragraphs[0].font.name = font
                        table.cell(row+1, col).text_frame.paragraphs[0].font.bold = bold
                        table.cell(row+1, col).text_frame.paragraphs[0].font.color.rgb = RGBColor(font_red, font_green,
                                                                                                font_blue)
                        table.cell(row+1, col).fill.solid()
                        table.cell(row+1, col).fill.fore_color.rgb = RGBColor(cell_red, cell_green, cell_blue)


def df_to_table_AMM_layout_v2(presentation, slidenumber, shape_name,
                            dataframe,table_header,
                            font='Amazon Ember Light', font_size=12, bold=False,
                            alignment='CENTER',
                            font_red=120, font_green=120, font_blue=120,color_code = MSO_THEME_COLOR.BACKGROUND_2):

    from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
    alignment_dictonary = {
            'CENTER': PP_PARAGRAPH_ALIGNMENT.CENTER,
            'LEFT': PP_PARAGRAPH_ALIGNMENT.LEFT,
            'RIGHT': PP_PARAGRAPH_ALIGNMENT.RIGHT,
            'JUSTIFY': PP_PARAGRAPH_ALIGNMENT.JUSTIFY}
    slide = presentation.slides[slidenumber - 1]
    from pptx.util import Inches
    x, y, cx, cy = Inches(4), Inches(3), Inches(20), Inches(0)
            
    shape=slide.shapes.add_table(dataframe.shape[0]+1, 2,x, y, cx, cy)
    table = shape.table    
    
    # Setting up merged header
    cell_00 = table.cell(0,0)
    cell_01 = table.cell(0,1)
    cell_00.merge(cell_01)
    
    cell_00.text = table_header   
    cell_00.text_frame.paragraphs[0].alignment = alignment_dictonary[alignment]
    cell_00.text_frame.paragraphs[0].font.size = Pt(font_size + 10)
    cell_00.text_frame.paragraphs[0].font.name = font
    cell_00.text_frame.paragraphs[0].font.bold = bold
    cell_00.text_frame.paragraphs[0].font.color.theme_color = color_code
    cell_00.fill.solid()
    cell_00.fill.fore_color.rgb = RGBColor(255, 255, 255)
  
    for row in range(dataframe.shape[0]):
                if (row % 2) == 0:
                    cell_red = 242
                    cell_green = 242
                    cell_blue = 242

                else:
                    cell_red = 255
                    cell_green = 255
                    cell_blue = 255

                #print('row: ', row)
                for col in range(dataframe.shape[1]):
                        # cell start with 0, 0 as top left column
                        # starting at row+1 we start from the first row in the PPT table
                      
                        cell = table.cell(row+1, col)  # row, columns
                        cell.text_frame.clear()
                        cell.text = str(dataframe.values[row][col])
                            
                        cell.text_frame.paragraphs[0].alignment = alignment_dictonary[alignment]
                        table.cell(row+1, col).text_frame.paragraphs[0].font.size = Pt(font_size)
                        table.cell(row+1, col).text_frame.paragraphs[0].font.name = font
                        table.cell(row+1, col).text_frame.paragraphs[0].font.bold = bold
                        table.cell(row+1, col).text_frame.paragraphs[0].font.color.rgb = RGBColor(font_red, font_green,
                                                                                                font_blue)
                        table.cell(row+1, col).fill.solid()
                        table.cell(row+1, col).fill.fore_color.rgb = RGBColor(cell_red, cell_green, cell_blue)
